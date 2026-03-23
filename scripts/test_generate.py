"""
Tech-Slides-Generator - Local Test Script
Created by: AresSheng (AI Product Manager, China)
License: MIT
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def generate_sample():
    print("🚀 AresSheng's Tech-Slides-Generator: Generating sample...")
    
    # 路径处理：确保能找到 assets 里的模板
    base_dir = os.path.dirname(os.path.dirname(__file__))
    template_path = os.path.join(base_dir, 'assets', 'base-template.pptx')
    
    if os.path.exists(template_path):
        prs = Presentation(template_path)
        print("✅ Base template loaded successfully.")
    else:
        prs = Presentation()
        print("⚠️ Template not found. Creating a default 4:3 presentation.")

    # 添加一张幻灯片
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # 使用空白版式

    # 1. 设置背景色 (Cyber-Blue: 5, 10, 31)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(5, 10, 31)

    # 2. 添加装饰性线条 (Accent: 0, 212, 255)
    line = slide.shapes.add_connector(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(9.5), Inches(1.2)
    )
    # 注意：pptx设置线条颜色较复杂，此处简化演示
    
    # 3. 添加主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "TECH SLIDES GENERATOR"
    run.font.bold = True
    run.font.size = Pt(44)
    run.font.color.rgb = RGBColor(255, 255, 255)

    # 4. 添加一个 Bento Box 风格的半透明面板
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2), Inches(5), Inches(4)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(15, 25, 60)
    panel.line.color.rgb = RGBColor(0, 212, 255)
    panel.line.width = Pt(1.5)

    # 保存
    output_path = "AresSheng_Sample_Output.pptx"
    prs.save(output_path)
    print(f"🎉 Success! Sample slide saved as: {output_path}")

if __name__ == "__main__":
    generate_sample()